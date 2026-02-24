#!/usr/bin/env python3
"""
Create a PowerPoint presentation from slide images.
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from PIL import Image

# Configuration
SLIDES_DIR = r"C:\Projects\DavinciMCP\output\slides"
OUTPUT_FILE = r"C:\Projects\DavinciMCP\output\dating_video_presentation.pptx"

# 16:9 dimensions: 1920x1080 = 13.333" x 7.5"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Create presentation with blank slide layout
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Get all slide PNGs sorted by name
slides_dir_path = Path(SLIDES_DIR)
slide_images = sorted([f for f in slides_dir_path.glob("slide_*.png")])

print(f"Found {len(slide_images)} slide images")

for idx, image_path in enumerate(slide_images, 1):
    print(f"Processing {idx}/31: {image_path.name}")

    # Create blank slide (use blank layout)
    blank_slide_layout = prs.slide_layouts[6]  # 6 is typically blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # Get image dimensions to scale properly
    with Image.open(image_path) as img:
        img_width, img_height = img.size

    # Calculate scale to fit image to slide (maintain aspect ratio, fill slide)
    slide_aspect = (13.333 / 7.5)
    image_aspect = img_width / img_height

    if image_aspect > slide_aspect:
        # Image is wider than slide aspect
        new_width = SLIDE_WIDTH
        new_height = int(SLIDE_WIDTH * img_height / img_width)
    else:
        # Image is taller than slide aspect
        new_height = SLIDE_HEIGHT
        new_width = int(SLIDE_HEIGHT * img_width / img_height)

    # Add image to slide at (0, 0) to cover entire slide
    left = Inches(0)
    top = Inches(0)

    # Add picture
    slide.shapes.add_picture(str(image_path), left, top, width=new_width, height=new_height)

# Save presentation
prs.save(OUTPUT_FILE)
print(f"\nPresentation saved to {OUTPUT_FILE}")

# Check file size
if os.path.exists(OUTPUT_FILE):
    file_size = os.path.getsize(OUTPUT_FILE)
    file_size_mb = file_size / (1024 * 1024)
    print(f"Output file size: {file_size:,} bytes ({file_size_mb:.2f} MB)")
else:
    print("ERROR: Output file not found!")
